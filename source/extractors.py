from pptx import Presentation
from docx import Document
import os
import fitz # PyMuPDF

# Extract headlines and paragraphs from DOCX file
def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        headlines = [para.text for para in doc.paragraphs if para.style.name.startswith('Heading')]
        return {'headlines': headlines, 'paragraphs': paragraphs}
    except Exception as _:
        return {'headlines': [], 'paragraphs': []}

# Extract headlines, paragraphs, and images from PPTX file
# Extract headlines, paragraphs, and images from PPTX file
def extract_content_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        headlines, paragraphs, images = [], [], []
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # Extract text
                if hasattr(shape, "text"):
                    if shape.text_frame:
                        paragraphs.append(shape.text)
                
                # Extract images
                if hasattr(shape, "image"):
                    image = shape.image
                    image_bytes = image.blob  # Get the image binary data
                    image_extension = os.path.splitext(image.content_type)[1]  # Get the image extension

                    # Create image path to save
                    image_path = os.path.join(os.getcwd(), 'static', 'uploads', f'image_slide{slide_idx}.{image_extension}')
                    
                    # Write image data to file
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                        images.append(image_path)

        return {'headlines': headlines, 'paragraphs': paragraphs, 'images': images}
    except Exception as _:
        return {'headlines': [], 'paragraphs': [], 'images': []}

def extract_content_from_pdf(file_path):
    try:
        doc = fitz.open(file_path)
        paragraphs, images = [], []

        for page_idx, page in enumerate(doc):
            # Extract text
            text = page.get_text("text").strip()
            if text:
                paragraphs.append(text)
            
            # Extract images
            for img_idx, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]  # Image extension
                
                # Create image path to save
                image_path = os.path.join(os.getcwd(), 'static', 'uploads', f'pdf_page{page_idx}_img{img_idx}.{image_ext}')
                
                # Write image data to file
                with open(image_path, 'wb') as img_file:
                    img_file.write(image_bytes)
                    images.append(image_path)
        
        return {'headlines': paragraphs, 'paragraphs': paragraphs, 'images': images}
    except Exception as _:
        return {'headlines': [], 'paragraphs': [], 'images': []}

# More extractors can be added here (e.g., for DOC format using third-party libraries)
